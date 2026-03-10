from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Function to add a slide with title, content, and notes
def add_slide(prs, title, content_items, notes_text):
    # Choose layout: 0 = Title Slide, 1 = Title and Content
    layout_index = 0 if title == "Title Slide" else 1
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    if title == "Title Slide":
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = "Beyond Shyness:\nUnderstanding Social Anxiety Disorder"
        subtitle_shape.text = "Causes, Symptoms, and Paths to Recovery\n\n[Your Name]\n[Date]"
    else:
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set Content (Bullet points)
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_items[0] # First bullet
        
        for item in content_items[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Add Speaker Notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

# --- Slide Data ---

slides_content = [
    {
        "title": "Title Slide",
        "content": [],
        "notes": "Welcome everyone. Today we are going to look beyond the surface of 'shyness' to understand Social Anxiety Disorder."
    },
    {
        "title": "The 'Spotlight' Effect",
        "content": [
            "Scenario: Walking into a room where everyone stops and looks at you.",
            "The Feeling: Rapid heartbeat, sweating, urge to run.",
            "The Reality: Usually a passing moment.",
            "The Disorder: For those with Social Anxiety, this is daily life.",
            "Goal: To define, understand, and discuss solutions."
        ],
        "notes": "Start with a hook. Ask the audience if they’ve ever felt nervous before a speech. Explain that while nervousness is normal, Social Anxiety is a constant, intense fear of being watched that disrupts life."
    },
    {
        "title": "What is Social Anxiety Disorder (SAD)?",
        "content": [
            "Definition: An intense, persistent fear of being watched and judged.",
            "Key Distinction: It is NOT just 'shyness' or 'introversion'.",
            "The Core Fear: Humiliation, embarrassment, or rejection.",
            "Impact: Disrupts daily functioning and relationships."
        ],
        "notes": "Define the disorder formally. Emphasize that shyness is a personality trait, whereas Social Anxiety is a mental health condition."
    },
    {
        "title": "Who Does It Affect?",
        "content": [
            "Statistics: Affects approx. 7% of the population.",
            "Onset: Usually starts during teenage years (avg. age 13).",
            "Gender: Slightly more common in females.",
            "Treatment Seeking: Men often seek help sooner due to career impact."
        ],
        "notes": "Highlight that this is one of the most common mental health disorders. If you have it, you are not alone."
    },
    {
        "title": "Physical Symptoms",
        "content": [
            "Blushing",
            "Profuse sweating",
            "Trembling or shaking",
            "Rapid heart rate",
            "Nausea or 'butterflies'",
            "Mind going blank"
        ],
        "notes": "Explain the 'Fight or Flight' response. The body perceives a social situation as a life-threatening danger, releasing adrenaline."
    },
    {
        "title": "Psychological & Emotional Symptoms",
        "content": [
            "Intense worry days or weeks before an event.",
            "Fear that others will notice you look anxious.",
            "Fear of being judged as stupid, awkward, or boring.",
            "Expecting the worst-case scenario."
        ],
        "notes": "Discuss the internal monologue. It’s a constant stream of self-criticism and 'what if' scenarios."
    },
    {
        "title": "Behavioral Symptoms",
        "content": [
            "Avoidance: Skipping parties, school, or work meetings.",
            "Safety Behaviors: Looking at phone to avoid eye contact.",
            "Silence: Staying quiet to avoid attention.",
            "Need for a 'Buddy': Unable to go places alone."
        ],
        "notes": "Explain that 'Avoidance' is the biggest maintainer of anxiety. The more you avoid, the scarier the situation becomes."
    },
    {
        "title": "Common Triggers",
        "content": [
            "Public speaking (most common)",
            "Meeting new people / Strangers",
            "Eating or drinking in public",
            "Making phone calls",
            "Using public restrooms"
        ],
        "notes": "Triggers vary from person to person. Some fear performing (speaking), while others fear interaction (small talk)."
    },
    {
        "title": "The Cycle of Social Anxiety",
        "content": [
            "1. Anticipatory Anxiety: Worrying about the event beforehand.",
            "2. The Situation: Enduring the event with high distress.",
            "3. Post-Event Processing (Rumination): Replaying the event in your head for days.",
            "Result: Increases fear for the next time."
        ],
        "notes": "Spend time on 'Post-Event Processing.' This is where people agonize over a small stutter or a joke that fell flat for days after it happened."
    },
    {
        "title": "Causes and Risk Factors",
        "content": [
            "Genetics: It can run in families.",
            "Brain Structure: Overactive Amygdala (the fear center).",
            "Environment: Past bullying or negative social experiences.",
            "Parenting: Controlling or overprotective parenting styles."
        ],
        "notes": "It is usually a combination of nature (biology) and nurture (environment), not just one thing."
    },
    {
        "title": "Impact on Daily Life",
        "content": [
            "Academic: Lower grades due to avoiding participation.",
            "Career: Turning down promotions to avoid presentations.",
            "Social: Difficulty forming friendships; isolation.",
            "Health: Risk of depression or substance abuse."
        ],
        "notes": "This isn't just about feeling nervous; it holds people back from achieving their potential and living the life they want."
    },
    {
        "title": "Myths vs. Facts",
        "content": [
            "Myth: 'It's just a phase.' -> Fact: Without treatment, it can last a lifetime.",
            "Myth: 'Socially anxious people don't like people.' -> Fact: They crave connection but fear rejection.",
            "Myth: 'Alcohol helps.' -> Fact: It is a temporary crutch that worsens anxiety long-term."
        ],
        "notes": "Debunking these myths is crucial for reducing stigma."
    },
    {
        "title": "Professional Treatment Options",
        "content": [
            "CBT (Cognitive Behavioral Therapy): The Gold Standard.",
            "Exposure Therapy: Gradually facing feared situations.",
            "Medication: SSRIs (antidepressants) or Beta-blockers.",
            "Group Therapy: Practicing social skills in a safe environment."
        ],
        "notes": "Emphasize that SAD is highly treatable. CBT helps re-wire the brain to stop seeing social situations as threats."
    },
    {
        "title": "Self-Help Strategies",
        "content": [
            "Challenge Negative Thoughts: 'Is there evidence for this fear?'",
            "Deep Breathing: Controls the physical fight-or-flight response.",
            "Shift Focus Outward: Focus on the room, not your internal feelings.",
            "Small Steps: Do one small scary thing a day."
        ],
        "notes": "Give the audience a practical tip: When you feel anxious, focus on the color of the walls or the texture of the chair to ground yourself."
    },
    {
        "title": "How to Help a Friend",
        "content": [
            "Listen without judgment.",
            "Don't force them into situations they aren't ready for.",
            "Praise small steps forward.",
            "Ask: 'How can I support you right now?'",
            "Be patient."
        ],
        "notes": "Telling someone to 'just calm down' or 'get over it' never works. Patience is key."
    },
    {
        "title": "Famous People with Social Anxiety",
        "content": [
            "[Insert Photos Here]",
            "Adele (Singer)",
            "Ryan Reynolds (Actor)",
            "Ed Sheeran (Singer)",
            "Message: Success is possible despite anxiety."
        ],
        "notes": "Show that even performers who are in the spotlight struggle with this. It does not define your capability."
    },
    {
        "title": "Conclusion & Questions",
        "content": [
            "Summary: Social Anxiety is a common, manageable condition.",
            "Key Takeaway: You are not your anxiety. Help is available.",
            "Resources: Anxiety & Depression Association of America (ADAA).",
            "Any Questions?"
        ],
        "notes": "End on a high, encouraging note. Remind the audience that bravery isn't the absence of fear, but acting in spite of it."
    }
]

# Generate Slides
for slide_data in slides_content:
    add_slide(prs, slide_data["title"], slide_data["content"], slide_data["notes"])

# Save the presentation
file_path = "Social_Anxiety_Presentation.pptx"
prs.save(file_path)